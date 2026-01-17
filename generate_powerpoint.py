#!/usr/bin/env python3
"""
Generate PowerPoint presentation arguing against vector search for agentic use cases.

Requirements:
    pip install python-pptx

Usage:
    python generate_powerpoint.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define colors
    DARK_BLUE = RGBColor(44, 62, 80)
    BRIGHT_BLUE = RGBColor(52, 152, 219)
    RED = RGBColor(231, 76, 60)
    GREEN = RGBColor(39, 174, 96)
    GRAY = RGBColor(127, 140, 141)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BRIGHT_BLUE

    title = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title.text_frame
    tf.text = "Vector Search: Unnecessary Cost"
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1))
    tf = subtitle.text_frame
    tf.text = "Why BM25 + Agentic Search Wins for B2B Enterprise"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(236, 240, 241)
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    title = slide.shapes.title
    title.text = "Executive Summary"

    left = Inches(1)
    top = Inches(1.8)
    width = Inches(8)
    height = Inches(4.5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Recommendation: Use BM25-Only Search"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_after = Pt(12)

    p = tf.add_paragraph()
    p.text = "Vector/semantic search provides ZERO measurable quality improvement over BM25 for technical documentation, while adding:"
    p.font.size = Pt(18)
    p.space_after = Pt(12)

    for item in ["Higher costs (vs. free)", "50-130x worse latency",
                 "External API dependencies", "Operational complexity"]:
        p = tf.add_paragraph()
        p.text = "✗ " + item
        p.font.size = Pt(20)
        p.font.color.rgb = RED
        p.level = 0

    # Slide 3: Real-World Test
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Real-World Test: Identical Results"

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = 'Query: "Describe differences between Windows and Linux TCP/IP stacks"'
    p.font.size = Pt(18)
    p.font.italic = True

    # BM25 box
    left = Inches(1)
    top = Inches(2.7)
    width = Inches(3.8)
    height = Inches(3)

    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(212, 237, 218)
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(3)

    tf = shape.text_frame
    tf.margin_top = Inches(0.2)
    tf.margin_left = Inches(0.2)

    p = tf.paragraphs[0]
    p.text = "BM25-Only"
    p.font.size = Pt(24)
    p.font.bold = True

    for item in ["✓ Top result: Correct", "Latency: <10ms", "Cost: $0"]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        if item.startswith("✓"):
            p.font.color.rgb = GREEN

    # Vector box
    left = Inches(5.2)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 215, 218)
    shape.line.color.rgb = RED
    shape.line.width = Pt(3)

    tf = shape.text_frame
    tf.margin_top = Inches(0.2)
    tf.margin_left = Inches(0.2)

    p = tf.paragraphs[0]
    p.text = "Vector Hybrid"
    p.font.size = Pt(24)
    p.font.bold = True

    for item in ["✓ Same answer", "Latency: 500-1,300ms", "Cost: $2/year + API"]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        if "500-1,300ms" in item or "Cost:" in item:
            p.font.color.rgb = RED
        elif item.startswith("✓"):
            p.font.color.rgb = GREEN

    # Slide 4: Cost Breakdown
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Cost Breakdown: What Vector Search Actually Costs"

    # Create table
    rows, cols = 6, 3
    left = Inches(1.5)
    top = Inches(2)
    width = Inches(7)
    height = Inches(4)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(2)
    table.columns[2].width = Inches(2)

    # Header row
    headers = ["Cost Component", "BM25", "Vector Hybrid"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRIGHT_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    data = [
        ["Initial Indexing", "$0", "$0.013"],
        ["Per Query Cost", "$0", "$0.000002-0.000006"],
        ["Annual Cost (1K queries/day)", "$0", "~$2/year"],
        ["Query Latency", "<10ms", "500-1,300ms"],
        ["Storage", "267KB", "13MB + 267KB"]
    ]

    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)

            if col_idx == 1:  # BM25 column
                p.font.color.rgb = GREEN
                p.font.bold = True
            elif col_idx == 2:  # Vector column
                p.font.color.rgb = RED

    # Slide 5: What is Agentic Search
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Understanding Agentic Search"

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    tf = textbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Traditional Search: "
    p.font.size = Pt(18)
    p.font.bold = True
    run = p.add_run()
    run.text = "End-user types query → System searches → Results"
    run.font.bold = False
    p.space_after = Pt(12)

    p = tf.add_paragraph()
    p.text = "Agentic Search: "
    p.font.size = Pt(18)
    p.font.bold = True
    run = p.add_run()
    run.text = "User asks LLM → "
    run.font.bold = False
    run = p.add_run()
    run.text = "LLM formulates optimized search queries"
    run.font.color.rgb = BRIGHT_BLUE
    run.font.bold = True
    run = p.add_run()
    run.text = " → System searches → LLM synthesizes answer"
    run.font.bold = False
    p.space_after = Pt(20)

    p = tf.add_paragraph()
    p.text = "The Key Insight"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BRIGHT_BLUE
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "The LLM already performs semantic-to-keyword translation!"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_after = Pt(12)

    p = tf.add_paragraph()
    p.text = 'User: "How do Windows and Linux handle TCP/IP differently?"'
    p.font.size = Pt(16)
    p.font.italic = True

    p = tf.add_paragraph()
    p.text = 'LLM: "Windows Linux TCP socket differences Winsock closesocket fork CreateProcess"'
    p.font.size = Pt(16)
    p.font.color.rgb = GREEN
    p.font.bold = True
    p.space_after = Pt(12)

    p = tf.add_paragraph()
    p.text = "This is already the perfect BM25 query. Vector search adds nothing."
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED

    # Slide 6: Why LLMs Make Vector Search Redundant
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Why LLMs Make Vector Search Redundant"

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    tf = textbox.text_frame

    p = tf.paragraphs[0]
    p.text = "What LLMs Do Automatically:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.space_after = Pt(12)

    items = [
        "Keyword extraction - Identify salient technical terms",
        "Terminology precision - Use exact API names, not paraphrases",
        "Query reformulation - Retry with different terms if needed",
        "Multi-query strategies - Issue parallel queries from different angles"
    ]

    for item in items:
        p = tf.add_paragraph()
        p.text = "✓ " + item
        p.font.size = Pt(18)
        p.level = 0
        p.space_after = Pt(6)

    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(12)

    p = tf.add_paragraph()
    p.text = "Vector search bridges semantic gaps:"
    p.font.size = Pt(18)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = '"car repair" → "automotive maintenance"'
    p.font.size = Pt(16)
    p.font.italic = True

    p = tf.add_paragraph()
    p.text = "BUT: With LLM agents, the LLM already bridges these gaps before calling search!"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED

    # Slide 7: Why BM25 Wins
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Why BM25 Dominates for Technical Documentation"

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    tf = textbox.text_frame

    p = tf.paragraphs[0]
    p.text = "1. Precise Terminology Matching"
    p.font.size = Pt(22)
    p.font.bold = True
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "Technical docs use exact terms: Winsock, closesocket, WSAStartup, fork(), POSIX"
    p.font.size = Pt(16)
    p.space_after = Pt(6)

    p = tf.add_paragraph()
    p.text = "These are not synonyms - they're literal API names"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(16)

    p = tf.add_paragraph()
    p.text = "2. No Semantic Ambiguity"
    p.font.size = Pt(22)
    p.font.bold = True
    p.space_after = Pt(8)

    items = [
        "fork() doesn't mean 'split' - it's a system call",
        "socket has one technical meaning",
        "Precision matters more than semantic similarity"
    ]

    for item in items:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.level = 0

    # Slide 8: Performance Table
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Head-to-Head Performance"

    rows, cols = 7, 4
    left = Inches(0.8)
    top = Inches(1.8)
    width = Inches(8.4)
    height = Inches(4.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(2.2)
    table.columns[3].width = Inches(1.8)

    # Headers
    headers = ["Metric", "BM25-Only", "Hybrid (BM25+Vector)", "Winner"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRIGHT_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # Data
    data = [
        ["Top-1 Accuracy", "✓ Correct", "✓ Correct", "Tie"],
        ["Top-5 Relevance", "High", "High", "Tie"],
        ["Query Latency", "<10ms", "500-1300ms", "BM25"],
        ["Cost per Query", "$0", "~$0.000006", "BM25"],
        ["Complexity", "Low", "Medium", "BM25"],
        ["Failure Modes", "Disk I/O", "Disk I/O + API", "BM25"]
    ]

    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)

            if col_idx == 3 and "BM25" in cell_text:
                p.font.color.rgb = GREEN
                p.font.bold = True

    # Slide 9: When You Actually Need Vectors
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "When Vector Search Actually Adds Value"

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    tf = textbox.text_frame

    p = tf.paragraphs[0]
    p.text = "Vector search provides benefits when:"
    p.font.size = Pt(20)
    p.space_after = Pt(12)

    items = [
        "Natural language from end-users (not LLM-mediated)",
        "Inconsistent terminology across documents",
        "Cross-lingual search requirements",
        "Conceptual similarity over keyword matching",
        "Large, diverse corpus (1000s+ documents with varied terminology)"
    ]

    for item in items:
        p = tf.add_paragraph()
        p.text = "✓ " + item
        p.font.size = Pt(16)
        p.level = 0

    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(20)

    p = tf.add_paragraph()
    p.text = "None of these apply to agentic search over technical reference documentation."
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RED

    # Slide 10: The Bottom Line
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "The Bottom Line"
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    textbox = slide.shapes.add_textbox(Inches(2), Inches(2.8), Inches(6), Inches(1))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "∞ ROI"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = BRIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    textbox = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.7))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "BM25 delivers equivalent quality at zero cost"
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER

    # Box
    left = Inches(1.5)
    top = Inches(5)
    width = Inches(7)
    height = Inches(1.5)

    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(212, 237, 218)
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(3)

    tf = shape.text_frame
    tf.margin_top = Inches(0.2)

    p = tf.paragraphs[0]
    p.text = "For Agentic Search Over Technical Documentation:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Vector search is an unnecessary cost that provides zero value."
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RED
    p.alignment = PP_ALIGN.CENTER

    # Slide 11: Recommendations
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Recommendations"

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    tf = textbox.text_frame

    p = tf.paragraphs[0]
    p.text = "Immediate Actions:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_after = Pt(12)

    p = tf.add_paragraph()
    p.text = "1. Deploy BM25-only search for the MCP server"
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "2. Monitor query performance metrics"
    p.font.size = Pt(18)

    subitems = ["Query success rate", "Result relevance (spot-checks)", "Failed queries"]
    for item in subitems:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.level = 1

    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(12)

    p = tf.add_paragraph()
    p.text = "Reconsider If/When:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.space_after = Pt(12)

    items = [
        ">5% query failure rate on relevant content",
        "User queries use synonyms not in documentation",
        "Corpus grows to 1000+ documents with varied terminology"
    ]

    for item in items:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.level = 0

    # Slide 12: Questions
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BLUE

    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    textbox = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(2))
    tf = textbox.text_frame

    p = tf.paragraphs[0]
    p.text = "Key Takeaway:"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(236, 240, 241)
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(12)

    p = tf.add_paragraph()
    p.text = "BM25 + LLM Agents = Free, Fast, and Effective"
    p.font.size = Pt(26)
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "Vector Search = Expensive, Slow, and Redundant"
    p.font.size = Pt(26)
    p.font.color.rgb = RED
    p.alignment = PP_ALIGN.CENTER

    return prs

if __name__ == "__main__":
    print("Generating PowerPoint presentation...")
    prs = create_presentation()
    filename = "vector-search-unnecessary-cost.pptx"
    prs.save(filename)
    print(f"✓ Presentation saved as: {filename}")
    print("\nPresentation contains 12 slides:")
    print("  1. Title slide")
    print("  2. Executive Summary")
    print("  3. Real-World Test Results")
    print("  4. Cost Breakdown")
    print("  5. What is Agentic Search")
    print("  6. Why LLMs Make Vector Search Redundant")
    print("  7. Why BM25 Wins for Technical Docs")
    print("  8. Head-to-Head Performance")
    print("  9. When You Actually Need Vectors")
    print(" 10. The Bottom Line")
    print(" 11. Recommendations")
    print(" 12. Questions/Closing")
